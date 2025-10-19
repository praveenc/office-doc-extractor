#!/usr/bin/env python3
# /// script
# requires-python = ">=3.12"
# dependencies = [
#   "pypandoc>=1.13",
#   "pandas>=2.3.3",
#   "openpyxl>=3.1.5",
#   "python-pptx>=1.0.2",
#   "pdfplumber>=0.11.7",
#   "loguru>=0.7.3",
# ]
# ///
"""
Office Document Text Extractor.

Extracts text from Office documents (DOCX, XLSX, PPTX, PDF) and HTML files.
Supports: DOCX, XLSX, PPTX, PDF, HTML formats.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Final

import pandas as pd
import pypandoc
from loguru import logger

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PPTX support disabled.")

try:
    import pdfplumber

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("pdfplumber not installed. PDF support disabled.")
    logger.warning("python-pptx not installed. PPTX support disabled.")


class DocumentConverter:
    """Convert various document formats to PDF and extract text."""

    # Formats that pypandoc can handle directly
    PANDOC_FORMATS: Final[set[str]] = {".html", ".htm", ".docx", ".doc"}

    # Formats requiring special handling
    EXCEL_FORMATS: Final[set[str]] = {".xls", ".xlsx"}
    POWERPOINT_FORMATS: Final[set[str]] = {".ppt", ".pptx"}
    PDF_FORMATS: Final[set[str]] = {".pdf"}

    ALL_SUPPORTED: Final[set[str]] = (
        PANDOC_FORMATS | EXCEL_FORMATS | POWERPOINT_FORMATS | PDF_FORMATS
    )

    def __init__(self, pdf_engine: str = "pdflatex") -> None:
        """
        Initialize document converter.

        Args:
            pdf_engine: PDF engine to use (pdflatex, xelatex, lualatex)

        Raises:
            RuntimeError: If pandoc is not installed

        """
        self.pdf_engine = pdf_engine
        self._verify_pandoc()
        logger.info(f"Initialized converter with PDF engine: {pdf_engine}")

    def _verify_pandoc(self) -> None:
        """Verify pandoc is installed, download if needed."""
        try:
            version = pypandoc.get_pandoc_version()
            logger.info(f"Pandoc version: {version}")
        except OSError:
            # Try to download pandoc automatically
            logger.info("Pandoc not found. Attempting to download...")
            try:
                pypandoc.download_pandoc()
                version = pypandoc.get_pandoc_version()
                logger.success(f"Pandoc downloaded successfully. Version: {version}")
            except Exception as e:
                msg = (
                    "Pandoc not found and auto-download failed. Install manually with:\n"
                    "  macOS: brew install pandoc\n"
                    "  Ubuntu: sudo apt-get install pandoc\n"
                    f"  Or run: python -m pypandoc.download_pandoc\n"
                    f"Error: {e}"
                )
                raise RuntimeError(msg) from e

    def validate_file(self, file_path: str) -> Path:
        """
        Validate input file.

        Args:
            file_path: Path to input file

        Returns:
            Validated Path object

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If format is not supported

        """
        path = Path(file_path)

        if not path.exists():
            logger.error(f"File not found: {file_path}")
            msg = f"File not found: {file_path}"
            raise FileNotFoundError(msg)

        if not path.is_file():
            logger.error(f"Path is not a file: {file_path}")
            msg = f"Path is not a file: {file_path}"
            raise ValueError(msg)

        suffix = path.suffix.lower()
        if suffix not in self.ALL_SUPPORTED:
            logger.error(f"Unsupported format: {suffix}")
            supported = ", ".join(sorted(self.ALL_SUPPORTED))
            msg = f"Unsupported format: {suffix}. Supported: {supported}"
            raise ValueError(msg)

        logger.info(f"Validated file: {path.name}")
        return path

    def convert_to_pdf(
        self,
        input_path: str,
        output_path: str | None = None,
        output_dir: str | None = None,
    ) -> Path:
        """
        Convert document to PDF.

        Args:
            input_path: Path to input document
            output_path: Path for output PDF (optional, auto-generated if None)
            output_dir: Directory for output PDF (optional, overrides output_path)

        Returns:
            Path to generated PDF file

        Raises:
            FileNotFoundError: If input file doesn't exist
            ValueError: If format is not supported
            RuntimeError: If conversion fails

        """
        input_file = self.validate_file(input_path)

        if output_dir is not None:
            output_directory = Path(output_dir)
            output_directory.mkdir(parents=True, exist_ok=True)
            output_file = output_directory / input_file.with_suffix(".pdf").name
        elif output_path is not None:
            output_file = Path(output_path)
        else:
            output_file = input_file.with_suffix(".pdf")

        suffix = input_file.suffix.lower()

        try:
            if suffix in self.PANDOC_FORMATS:
                self._convert_with_pandoc(input_file, output_file)
            elif suffix in self.EXCEL_FORMATS:
                self._convert_excel(input_file, output_file)
            elif suffix in self.POWERPOINT_FORMATS:
                self._convert_powerpoint(input_file, output_file)
            else:
                msg = f"Unsupported format: {suffix}"
                raise ValueError(msg)

            logger.success(f"Converted {input_file.name} -> {output_file.name}")
            return output_file

        except Exception as e:
            logger.error(f"Conversion failed: {e}")
            raise

    def extract_text(self, input_path: str) -> str:
        """
        Extract text directly from document without PDF conversion.

        This is more efficient for supported formats (HTML, DOCX) as it
        bypasses PDF conversion and Textract API calls.

        Args:
            input_path: Path to input document

        Returns:
            Extracted text content

        Raises:
            FileNotFoundError: If input file doesn't exist
            ValueError: If format is not supported
            RuntimeError: If extraction fails

        """
        input_file = self.validate_file(input_path)
        suffix = input_file.suffix.lower()

        try:
            if suffix in self.PANDOC_FORMATS:
                logger.info(f"Extracting text from {input_file.name} with pypandoc")
                text = pypandoc.convert_file(str(input_file), "plain")
                logger.success(f"Extracted {len(text)} characters")
                return text
            if suffix in self.EXCEL_FORMATS:
                logger.info(f"Extracting text from Excel file: {input_file.name}")
                return self._extract_excel_text(input_file)
            if suffix in self.POWERPOINT_FORMATS:
                logger.info(f"Extracting text from PowerPoint file: {input_file.name}")
                return self._extract_powerpoint_text(input_file)
            if suffix in self.PDF_FORMATS:
                logger.info(f"Extracting text from PDF file: {input_file.name}")
                return self._extract_pdf_text(input_file)
            msg = f"Unsupported format: {suffix}"
            raise ValueError(msg)

        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            raise

    def _convert_with_pandoc(self, input_file: Path, output_file: Path) -> None:
        """
        Convert document using pypandoc.

        Args:
            input_file: Input file path
            output_file: Output PDF path

        """
        logger.info(f"Converting {input_file.name} with pypandoc")

        try:
            pypandoc.convert_file(
                str(input_file),
                "pdf",
                outputfile=str(output_file),
                extra_args=[f"--pdf-engine={self.pdf_engine}"],
            )
        except RuntimeError as e:
            error_msg = str(e)
            if "pdflatex not found" in error_msg or "pdf-engine" in error_msg:
                msg = (
                    "PDF engine not found. Install LaTeX:\n"
                    "  macOS: brew install --cask mactex\n"
                    "  Ubuntu: sudo apt-get install texlive-latex-base"
                )
                raise RuntimeError(msg) from e
            raise

    def _convert_excel(self, input_file: Path, output_file: Path) -> None:
        """
        Convert Excel file to PDF.

        Args:
            input_file: Input Excel file path
            output_file: Output PDF path

        """
        logger.info(f"Converting Excel file: {input_file.name}")

        try:
            excel_file = pd.ExcelFile(input_file)
            html_parts = []

            for sheet_name in excel_file.sheet_names:
                data_frame = pd.read_excel(excel_file, sheet_name=sheet_name)
                html_parts.append(f"<h2>{sheet_name}</h2>")
                html_parts.append(data_frame.to_html(index=False, border=1))

            full_html = "\n".join(html_parts)

            styled_html = f"""
            <html>
            <head>
                <style>
                    body {{ font-family: Arial, sans-serif; margin: 20px; }}
                    h2 {{ color: #333; margin-top: 30px; }}
                    table {{ border-collapse: collapse; width: 100%; margin-bottom: 20px; }}
                    th {{ background-color: #4CAF50; color: white; padding: 8px; text-align: left; }}
                    td {{ padding: 8px; border: 1px solid #ddd; }}
                    tr:nth-child(even) {{ background-color: #f2f2f2; }}
                </style>
            </head>
            <body>
                {full_html}
            </body>
            </html>
            """

            pypandoc.convert_text(
                styled_html,
                "pdf",
                format="html",
                outputfile=str(output_file),
                extra_args=[f"--pdf-engine={self.pdf_engine}"],
            )

        except Exception as e:
            msg = f"Failed to convert Excel file: {e}"
            raise RuntimeError(msg) from e

    def _extract_excel_text(self, input_file: Path) -> str:
        """Extract text from Excel file."""
        excel_file = pd.ExcelFile(input_file)
        text_parts = []

        for sheet_name in excel_file.sheet_names:
            data_frame = pd.read_excel(excel_file, sheet_name=sheet_name)
            text_parts.append(f"Sheet: {sheet_name}\n")
            text_parts.append(data_frame.to_string(index=False))
            text_parts.append("\n\n")

        return "\n".join(text_parts)

    def _convert_powerpoint(self, input_file: Path, output_file: Path) -> None:
        """
        Convert PowerPoint file to PDF.

        Args:
            input_file: Input PowerPoint file path
            output_file: Output PDF path

        """
        if not PPTX_AVAILABLE:
            msg = (
                "python-pptx not installed. Install with: pip install python-pptx\n"
                "Note: Only .pptx format is supported, not .ppt"
            )
            raise RuntimeError(msg)

        if input_file.suffix.lower() == ".ppt":
            msg = (
                "Legacy .ppt format not supported. "
                "Please convert to .pptx format first using PowerPoint or LibreOffice."
            )
            raise ValueError(msg)

        logger.info(f"Converting PowerPoint file: {input_file.name}")

        try:
            markdown_content = self._extract_powerpoint_text(input_file)

            pypandoc.convert_text(
                markdown_content,
                "pdf",
                format="markdown",
                outputfile=str(output_file),
                extra_args=[
                    f"--pdf-engine={self.pdf_engine}",
                    "--variable=geometry:margin=1in",
                ],
            )

        except Exception as e:
            msg = f"Failed to convert PowerPoint file: {e}"
            raise RuntimeError(msg) from e

    def _extract_powerpoint_text(self, input_file: Path) -> str:
        """Extract text from PowerPoint file."""
        if not PPTX_AVAILABLE:
            msg = "python-pptx not installed"
            raise RuntimeError(msg)

        prs = Presentation(str(input_file))
        text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"# Slide {i}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if len(text) < 100 and "\n" not in text:
                        text_parts.append(f"## {text}\n")
                    else:
                        text_parts.append(f"{text}\n")

            text_parts.append("\n---\n\n")

        return "\n".join(text_parts)

    def _extract_pdf_text(self, input_file: Path) -> str:
        """
        Extract text from PDF file using pdfplumber.

        Args:
            input_file: Input PDF file path

        Returns:
            Extracted text content

        Raises:
            RuntimeError: If pdfplumber is not installed or extraction fails

        """
        if not PDF_AVAILABLE:
            msg = (
                "pdfplumber not installed. Install with: pip install pdfplumber\n"
                "Note: This extracts text from digital PDFs. For scanned PDFs, use AWS Textract."
            )
            raise RuntimeError(msg)

        try:
            text_parts = []

            with pdfplumber.open(str(input_file)) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Processing {total_pages} pages from PDF")

                for page_num, page in enumerate(pdf.pages, 1):
                    # Extract text from page
                    page_text = page.extract_text()

                    if page_text:
                        text_parts.append(f"--- Page {page_num} ---\n")
                        text_parts.append(page_text)
                        text_parts.append("\n\n")
                    else:
                        logger.warning(f"No text found on page {page_num}")

            full_text = "".join(text_parts)

            if not full_text.strip():
                logger.warning(
                    "No text extracted from PDF. This may be a scanned document.",
                )
                msg = (
                    "No text found in PDF. This appears to be a scanned document. "
                    "Use AWS Textract for OCR processing."
                )
                raise ValueError(msg)

            logger.success(
                f"Extracted {len(full_text)} characters from {total_pages} pages",
            )
            return full_text

        except Exception as e:
            if "No text found" in str(e):
                raise
            msg = f"Failed to extract text from PDF: {e}"
            raise RuntimeError(msg) from e


def main() -> int:
    """Main function to handle command line arguments and execute conversion."""
    parser = argparse.ArgumentParser(
        description="Convert documents to PDF for Amazon Textract processing",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats: HTML, PDF, DOCX, XLS, XLSX, PPTX

Examples:
  python office_doc_extractor.py document.html
  python office_doc_extractor.py document.docx --output result.pdf
  python office_doc_extractor.py spreadsheet.xlsx --output-dir converted/
  python office_doc_extractor.py --pdf-engine xelatex document.html
  python office_doc_extractor.py --extract-text document.docx
        """,
    )

    parser.add_argument("input_file", help="Path to input document")
    parser.add_argument(
        "--output",
        "-o",
        help="Output file path (for single file conversion)",
    )
    parser.add_argument(
        "--output-dir",
        "-d",
        default="output",
        help="Base output directory (default: output/)",
    )
    parser.add_argument(
        "--extract-text",
        "-t",
        action="store_true",
        help="Extract text directly without PDF conversion (faster, no Textract needed)",
    )
    parser.add_argument(
        "--pdf-engine",
        default="pdflatex",
        choices=["pdflatex", "xelatex", "lualatex"],
        help="PDF engine to use (default: pdflatex)",
    )
    parser.add_argument(
        "--quiet",
        "-q",
        action="store_true",
        help="Suppress logging output",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose logging",
    )

    args = parser.parse_args()

    # Configure logging
    logger.remove()
    if not args.quiet:
        log_level = "DEBUG" if args.verbose else "INFO"
        logger.add(
            sys.stderr,
            format="<green>{time:YYYY-MM-DD HH:mm:ss}</green> | <level>{level: <8}</level> | <level>{message}</level>",
            level=log_level,
        )

    try:
        converter = DocumentConverter(pdf_engine=args.pdf_engine)
        input_path = Path(args.input_file)

        if args.extract_text:
            # Direct text extraction
            text = converter.extract_text(args.input_file)

            if args.output:
                # Save to specific output file
                output_file = Path(args.output)
                output_file.parent.mkdir(parents=True, exist_ok=True)
                output_file.write_text(text, encoding="utf-8")
                if not args.quiet:
                    print(f"✓ Text saved to: {output_file}")
            else:
                # Save to output_dir/text/ folder
                text_dir = Path(args.output_dir) / "text"
                text_dir.mkdir(parents=True, exist_ok=True)
                output_file = text_dir / f"{input_path.stem}.txt"
                output_file.write_text(text, encoding="utf-8")
                if not args.quiet:
                    print(f"✓ Text saved to: {output_file}")
        else:
            # Convert to PDF
            if args.output:
                # Use specific output path
                output_file = converter.convert_to_pdf(
                    args.input_file,
                    output_path=args.output,
                )
            else:
                # Save to output_dir/converted_pdfs/ folder
                pdf_dir = Path(args.output_dir) / "converted_pdfs"
                pdf_dir.mkdir(parents=True, exist_ok=True)
                output_file = converter.convert_to_pdf(
                    args.input_file,
                    output_dir=str(pdf_dir),
                )

            if not args.quiet:
                print(f"✓ Successfully converted to: {output_file}")

        return 0

    except (FileNotFoundError, ValueError, RuntimeError) as e:
        logger.error(str(e))
        print(f"Error: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
